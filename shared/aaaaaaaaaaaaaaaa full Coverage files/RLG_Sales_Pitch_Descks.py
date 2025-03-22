#!/usr/bin/env python3
"""
RLG Sales Pitch Deck Generator
----------------------------------------
This script generates a complete PowerPoint (.pptx) sales pitch deck for RLG Data and RLG Fans.

It covers:
- Introduction & Overview
- Competitive Advantage
- AI-Powered Features
- RLG Super Tool
- Pricing & Monetization
- Case Studies & Testimonials
- Call to Action

The generated PowerPoint is fully formatted and ready to use.
"""

from pptx import Presentation
from pptx.util import Inches

# Create a PowerPoint presentation object
prs = Presentation()

# Theme Colors
TITLE_COLOR = "003366"
SUBTITLE_COLOR = "4B8BBE"
TEXT_COLOR = "000000"

# Function to add a slide
def add_slide(title, content, layout=1):
    """Creates a PowerPoint slide with a title and content."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout])
    title_box = slide.shapes.title
    content_box = slide.placeholders[1]

    title_box.text = title
    content_box.text = content

# 1. Title Slide
title_slide = prs.slides.add_slide(prs.slide_layouts[0])
title_slide.shapes.title.text = "RLG Data & RLG Fans Sales Pitch"
title_slide.shapes.placeholders[1].text = "Empowering Businesses with Real-Time Media Intelligence"

# 2. Problem Statement
add_slide(
    "The Challenge",
    "In today's fast-paced digital landscape, businesses struggle to track brand mentions, analyze sentiment, and measure compliance across regions."
)

# 3. Our Solution
add_slide(
    "RLG Data & RLG Fans: The Ultimate AI-Powered Solution",
    "✓ Real-time media monitoring\n"
    "✓ AI-driven sentiment & trend analysis\n"
    "✓ Automated compliance tracking\n"
    "✓ Competitive intelligence insights\n"
    "✓ Region-specific accuracy (country, city, town level)"
)

# 4. Competitive Advantage
add_slide(
    "How We Compare to Brandwatch, Mention & Meltwater",
    "RLG Data & RLG Fans is:\n"
    "✓ 40% more accurate in regional media tracking\n"
    "✓ 30% more efficient in sentiment analysis\n"
    "✓ More affordable with flexible pricing\n"
    "✓ Fully automated & AI-driven\n"
    "✓ Seamlessly integrated with compliance monitoring"
)

# 5. AI-Powered Features
add_slide(
    "Key Features",
    "✓ Smart AI-powered scraping\n"
    "✓ Sentiment & trend detection\n"
    "✓ GDPR & CCPA compliance automation\n"
    "✓ Custom alerts & notifications\n"
    "✓ Multi-platform accessibility (web, Chrome, mobile)"
)

# 6. RLG Super Tool
add_slide(
    "RLG Super Tool: A Game Changer",
    "✓ One unified dashboard for all media insights\n"
    "✓ AI-driven predictive analytics\n"
    "✓ Integration with all social & news platforms\n"
    "✓ Custom reports & real-time alerts"
)

# 7. Pricing & Monetization
add_slide(
    "Flexible Pricing Plans",
    "✓ Weekly Plan: $15 (Global) / $35 (Special Region)\n"
    "✓ Monthly Plan: $59 (Global) / $99 (Special Region)\n"
    "✓ Custom Enterprise Packages Available"
)

# 8. Case Studies & Testimonials
add_slide(
    "Success Stories",
    "✓ Case Study 1: A global brand increased media coverage by 50% using RLG Data.\n"
    "✓ Case Study 2: A fintech company reduced compliance risks by 70%.\n"
    "✓ Case Study 3: A startup improved its online reputation tracking by 85%."
)

# 9. Call to Action
add_slide(
    "Join the Future of AI-Powered Media Monitoring",
    "✓ Sign up today and experience the power of RLG Data & RLG Fans!\n"
    "✓ Contact us at: support@rlgdata.com"
)

# Save the presentation
output_filename = "RLG_Sales_Pitch_Decks.pptx"
prs.save(output_filename)
print(f"Sales pitch deck generated: {output_filename}")
